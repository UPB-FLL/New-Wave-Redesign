from __future__ import annotations

import os
import zipfile
from pathlib import Path

import fitz
import pytest
from lxml import etree
from PIL import Image
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx import Presentation

from brandkit.cli import build
from brandkit.logo import export_logo_family
from brandkit.patterns import export_patterns
from brandkit.render import export_raster_family
from brandkit.slides import PRESENTATION_LAYOUTS, build_presentation_template


_RUNTIME_NODE = Path(
    "C:/Users/JoshBreault/.cache/codex-runtimes/codex-primary-runtime/dependencies/node/bin/node.exe"
)
_RUNTIME_MODULES = Path(
    "C:/Users/JoshBreault/.cache/codex-runtimes/codex-primary-runtime/dependencies/node/node_modules"
)
_RUNTIME_BIN = Path(
    "C:/Users/JoshBreault/.cache/codex-runtimes/codex-primary-runtime/dependencies/bin/override"
)
_PACKAGE_ROOT = "08-documents-presentations"
_PPTX_NAME = "new-wave-it-presentation-template.pptx"


@pytest.fixture(scope="module")
def package_root(tmp_path_factory: pytest.TempPathFactory) -> Path:
    if not (_RUNTIME_NODE.is_file() and _RUNTIME_MODULES.is_dir() and _RUNTIME_BIN.is_dir()):
        pytest.skip("The bundled artifact-tool runtime is unavailable in this environment.")

    os.environ.update(
        {
            "RUNTIME_NODE": str(_RUNTIME_NODE),
            "RUNTIME_NODE_MODULES": str(_RUNTIME_MODULES),
            "RUNTIME_BIN_DIR": str(_RUNTIME_BIN),
        }
    )
    root = build(tmp_path_factory.mktemp("slides") / "new-wave-it-brand-package")
    export_logo_family(root)
    export_raster_family(root)
    export_patterns(root)
    deck_root = root / _PACKAGE_ROOT
    (deck_root / "new-wave-it-presentation-template-debug.pptx").write_bytes(b"debug")
    (deck_root / "new-wave-it-presentation-template-debug.pptx.inspect.ndjson").write_text("debug", encoding="utf-8")
    build_presentation_template(deck_root / _PPTX_NAME)
    return root


def test_presentation_template_has_the_required_editable_layout_system(package_root: Path) -> None:
    target = package_root / _PACKAGE_ROOT / _PPTX_NAME
    assert target.is_file() and target.stat().st_size > 10_000

    presentation = Presentation(target)
    assert presentation.slide_width / presentation.slide_height == pytest.approx(16 / 9)
    assert set(PRESENTATION_LAYOUTS) <= {layout.name for layout in presentation.slide_layouts}
    assert len(presentation.slides) == len(PRESENTATION_LAYOUTS) + 1

    used_layouts = {slide.slide_layout.name for slide in list(presentation.slides)[:-1]}
    assert set(PRESENTATION_LAYOUTS) <= used_layouts
    for layout_name in PRESENTATION_LAYOUTS:
        layout = next(layout for layout in presentation.slide_layouts if layout.name == layout_name)
        placeholder_types = {placeholder.placeholder_format.type for placeholder in layout.placeholders}
        assert {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT} <= placeholder_types


def test_presentation_has_notes_hidden_reference_and_review_artifacts(package_root: Path) -> None:
    deck_root = package_root / _PACKAGE_ROOT
    review_pdf = deck_root / "new-wave-it-presentation-template-review.pdf"
    contact_sheet = deck_root / "new-wave-it-presentation-contact-sheet.png"
    assert review_pdf.is_file() and review_pdf.stat().st_size > 1_000
    assert contact_sheet.is_file() and contact_sheet.stat().st_size > 1_000
    assert not (package_root / ".brand-build").exists()
    assert not list(deck_root.glob("*debug*"))
    assert not list(deck_root.glob("*.inspect.ndjson"))

    review = fitz.open(review_pdf)
    try:
        assert review.page_count == len(PRESENTATION_LAYOUTS) + 1
        assert all(page.rect.width > page.rect.height for page in review)
    finally:
        review.close()

    with Image.open(contact_sheet) as image:
        assert image.width >= 960
        assert image.height >= 700

    with zipfile.ZipFile(deck_root / _PPTX_NAME) as archive:
        names = archive.namelist()
        notes = [name for name in names if name.startswith("ppt/notesSlides/notesSlide") and name.endswith(".xml")]
        assert len(notes) == len(PRESENTATION_LAYOUTS) + 1
        notes_text = "\n".join(archive.read(name).decode("utf-8", "ignore") for name in notes)
        assert "Editable template guidance" in notes_text
        reference = etree.fromstring(archive.read(f"ppt/slides/slide{len(PRESENTATION_LAYOUTS) + 1}.xml"))
        assert reference.get("show") in {"0", "false"}
        reference_text = " ".join(reference.xpath("//*[local-name()='t']/text()"))
        assert "INTERNAL REFERENCE" in reference_text
        assert "Title" in reference_text and "Closing" in reference_text


def test_presentation_uses_only_approved_color_tokens_and_no_gradients(package_root: Path) -> None:
    approved = {"09131D", "101E2D", "31C6CF", "62BE68", "317B92", "F7FAFB", "E8EFF1", "526271", "FFFFFF"}
    with zipfile.ZipFile(package_root / _PACKAGE_ROOT / _PPTX_NAME) as archive:
        xml_names = [name for name in archive.namelist() if name.endswith(".xml")]
        payload = "\n".join(archive.read(name).decode("utf-8", "ignore") for name in xml_names)
        roots = [etree.fromstring(archive.read(name)) for name in xml_names]

    direct_colors = {
        value.upper()
        for root in roots
        for value in root.xpath("//*[local-name()='srgbClr']/@val")
    }
    assert direct_colors <= approved
    assert "000000" not in direct_colors
    assert "gradFill" not in payload
